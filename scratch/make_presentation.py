import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Colors
BROWN = RGBColor(75, 56, 50)
GOLD = RGBColor(241, 210, 122)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(45, 30, 24)
ACCENT = RGBColor(163, 107, 79)

def add_header(slide, title_text, category_text="SWEET COUNTY BAKERY"):
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = header_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = DARK_TEXT

# SLIDE 1: Title Slide
slide1 = prs.slides.add_slide(blank_slide_layout)
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BROWN

tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "🍰 Sweet County Bakery"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = GOLD

p_sub = tf1.add_paragraph()
p_sub.text = "Full-Stack Artisanal E-Commerce Platform with 3D Spotlights & Animated Packing"
p_sub.font.size = Pt(20)
p_sub.font.color.rgb = WHITE

p_url = tf1.add_paragraph()
p_url.text = "\nLive Application: https://sweet-county.vercel.app"
p_url.font.size = Pt(16)
p_url.font.color.rgb = GOLD

# SLIDE 2: Architecture
slide2 = prs.slides.add_slide(blank_slide_layout)
add_header(slide2, "🏗️ System Architecture & Technology Stack")
tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf2 = tb2.text_frame
tf2.word_wrap = True

items2 = [
    ("Frontend Layer (React 19 / Vite)", "Built with React 19, Vanilla CSS with 3D Perspective, React Router v7, and Context API (Auth & Cart)."),
    ("Backend API Layer (Vercel Serverless)", "Node.js Express Serverless API handling authentication, catalog filtering, cart orders, and payments at /api/*."),
    ("Database & Storage", "Embedded In-Memory Data Store & MongoDB Atlas integration for zero-delay response times."),
    ("Deployment & Hosting", "Unified single-domain deployment on Vercel (https://sweet-county.vercel.app).")
]

for title, desc in items2:
    p_t = tf2.add_paragraph()
    p_t.text = "• " + title
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT
    
    p_d = tf2.add_paragraph()
    p_d.text = "   " + desc + "\n"
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = DARK_TEXT

# SLIDE 3: Features & 3D Interactive Design
slide3 = prs.slides.add_slide(blank_slide_layout)
add_header(slide3, "✨ Key Features & 3D User Experience")
tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf3 = tb3.text_frame
tf3.word_wrap = True

items3 = [
    ("3D Levitating Cake Spotlight", "Real-time mouse-tilt container with mid-air floating animation and orbital sparkles."),
    ("3D Bakery Box Packing Animation", "Cake drops into 3D Bakery Box -> Lid closes -> Ribbon wraps -> Wax seal stamps before payment."),
    ("Reassuring Checkout & Payment", "Razorpay test modal supporting UPI QR (GPay, PhonePe, Paytm), Cards, Netbanking & Cash on Delivery."),
    ("Admin Dashboard & User Database", "Live order pipeline management and registered customer account viewer.")
]

for title, desc in items3:
    p_t = tf3.add_paragraph()
    p_t.text = "• " + title
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT
    
    p_d = tf3.add_paragraph()
    p_d.text = "   " + desc + "\n"
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = DARK_TEXT

# SLIDE 4: Technical Challenges
slide4 = prs.slides.add_slide(blank_slide_layout)
add_header(slide4, "🚧 Technical Challenges & Solutions")
tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf4 = tb4.text_frame
tf4.word_wrap = True

items4 = [
    ("Challenge 1: Connection Errors ('Could not connect to server')", "Fix: Unified Frontend and Backend onto Vercel Serverless Functions (api/index.js)."),
    ("Challenge 2: 500 Serverless Execution Error", "Fix: Converted api/index.js to ES Module format (export default handler) matching package.json."),
    ("Challenge 3: 404 Error on Direct Page Load (/cart)", "Fix: Configured Single Page Application (SPA) catch-all rewrite in vercel.json to /index.html.")
]

for title, desc in items4:
    p_t = tf4.add_paragraph()
    p_t.text = "• " + title
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT
    
    p_d = tf4.add_paragraph()
    p_d.text = "   " + desc + "\n"
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = DARK_TEXT

os.makedirs("docs", exist_ok=True)
prs.save("docs/Sweet_County_Bakery_Presentation.pptx")
print("Presentation generated successfully at docs/Sweet_County_Bakery_Presentation.pptx!")
